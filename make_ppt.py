from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour Palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0A, 0x0F, 0x1E)   # deep navy
CARD_BG   = RGBColor(0x12, 0x1A, 0x35)   # slightly lighter navy
BLUE      = RGBColor(0x4F, 0x8C, 0xFF)   # accent blue
TEAL      = RGBColor(0x00, 0xD4, 0xAA)   # accent teal
PURPLE    = RGBColor(0xA8, 0x55, 0xF7)   # accent purple
PINK      = RGBColor(0xF4, 0x72, 0xB6)   # accent pink
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
LIGHT     = RGBColor(0xE2, 0xE8, 0xF0)

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper Functions ──────────────────────────────────────────────────────────

def add_bg(slide, color=DARK_BG):
    """Fill entire slide background."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_w=Pt(1), radius=False):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1=rectangle
    shape.line.width = Pt(0)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True,
             font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_multi_para(slide, lines, l, t, w, h,
                    font_size=13, color=LIGHT, bullet_char="◆ ", bold_first=False,
                    line_spacing_pt=None, font_name="Calibri"):
    """Add a textbox with multiple bullet lines."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet_char + line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
        if bold_first and i == 0:
            run.font.bold = True
    return txb

def add_header_bar(slide, title, subtitle="", accent=BLUE):
    """Add a coloured top bar with title and optional subtitle."""
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.55), fill_color=CARD_BG)
    # accent left stripe
    add_rect(slide, Inches(0), Inches(0), Inches(0.07), Inches(1.55), fill_color=accent)
    add_text(slide, title, Inches(0.3), Inches(0.18), Inches(12.5), Inches(0.8),
             font_size=30, bold=True, color=WHITE, font_name="Calibri")
    if subtitle:
        add_text(slide, subtitle, Inches(0.3), Inches(0.92), Inches(12.5), Inches(0.45),
                 font_size=13, color=accent, font_name="Calibri")

def add_card(slide, l, t, w, h, bg=CARD_BG, border=BLUE):
    add_rect(slide, l, t, w, h, fill_color=bg, line_color=border, line_w=Pt(1.2))

def card_title(slide, text, l, t, w, accent=TEAL, fs=13):
    add_text(slide, text, l, t, w, Inches(0.35), font_size=fs,
             bold=True, color=accent, font_name="Calibri")

def slide_number(slide, n, total=15):
    add_text(slide, f"{n} / {total}",
             Inches(11.8), Inches(7.1), Inches(1.3), Inches(0.3),
             font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — TITLE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)

# Gradient bar at top
add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.12), fill_color=BLUE)
# Bottom bar
add_rect(s, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3), fill_color=CARD_BG)

# Hospital emoji / logo box
add_rect(s, Inches(5.67), Inches(0.9), Inches(2), Inches(2), fill_color=BLUE)
add_text(s, "🏥", Inches(5.8), Inches(0.95), Inches(1.8), Inches(1.8),
         font_size=52, align=PP_ALIGN.CENTER)

# Main title
add_text(s, "HOSPITAL MANAGEMENT SYSTEM",
         Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.9),
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

add_text(s, "A Web-Based Enterprise Platform using Java Spring Boot, Spring Data JPA, H2 Database & SPA Dashboard",
         Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.6),
         font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Divider line
add_rect(s, Inches(3.5), Inches(4.6), Inches(6.3), Inches(0.04), fill_color=TEAL)

add_text(s, "Presented by: S Aditya  |  Enrollment: 22STUCRPN01019",
         Inches(0.5), Inches(4.75), Inches(12.3), Inches(0.4),
         font_size=13, color=TEAL, align=PP_ALIGN.CENTER, bold=True)
add_text(s, "M.Sc. Computer Science  |  ICFAI University, Raipur (C.G.)  |  Session: 2024–26",
         Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
         font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

slide_number(s, 1)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — TABLE OF CONTENTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Table of Contents", "What We Will Cover", accent=TEAL)

toc_items = [
    ("01", "Introduction & Problem Statement",    BLUE),
    ("02", "Objectives & Scope",                   TEAL),
    ("03", "Theoretical Background",               PURPLE),
    ("04", "System Architecture",                  PINK),
    ("05", "Patient & Doctor Management",           BLUE),
    ("06", "Appointment & Billing System",          TEAL),
    ("07", "REST API Design",                       PURPLE),
    ("08", "SPA Dashboard & UI",                    PINK),
    ("09", "Database Design (JPA / H2)",            BLUE),
    ("10", "Data Seeder & Testing",                 TEAL),
    ("11", "Security & Validation",                 PURPLE),
    ("12", "Deployment Strategy",                   PINK),
    ("13", "Real-World Clinical Evaluation",        BLUE),
    ("14", "Limitations & Future Scope",            TEAL),
    ("15", "Conclusion",                            PURPLE),
]

cols = [toc_items[:8], toc_items[8:]]
for ci, col in enumerate(cols):
    lx = Inches(0.35 + ci * 6.5)
    for ri, (num, title, clr) in enumerate(col):
        ty = Inches(1.75 + ri * 0.67)
        add_rect(s, lx, ty, Inches(0.45), Inches(0.44), fill_color=clr)
        add_text(s, num, lx, ty, Inches(0.45), Inches(0.44),
                 font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text(s, title, Inches(0.86 + ci * 6.5), ty + Pt(2),
                 Inches(5.8), Inches(0.44), font_size=12, color=LIGHT)

slide_number(s, 2)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — INTRODUCTION & PROBLEM STATEMENT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Introduction & Problem Statement", "Why does this system exist?", accent=PINK)

intro = (
    "Traditional healthcare facilities — especially localized clinics and smaller hospitals — "
    "suffer from fragmented and siloed operational workflows. Doctor availability scheduling, "
    "patient history logging, and billing generation are all handled in isolation, causing "
    "transcription errors, appointment conflicts, double-booking, and revenue loss."
)
add_text(s, intro, Inches(0.35), Inches(1.65), Inches(12.6), Inches(1.0),
         font_size=13, color=LIGHT)

problems = [
    ("📋  Information Silos",       "Ward nurses, doctors, and billing clerks work on separate systems (paper registers, prescription pads, manual ledgers) with no synchronization."),
    ("📅  Scheduling Conflicts",    "Appointments are booked without real-time visibility into doctor availability — resulting in overbooking, long queues, and clinician fatigue."),
    ("💸  Billing Omissions",       "Completed consultations must be manually reported to the billing counter. Patients leaving unnotified leads to direct revenue loss for the clinic."),
    ("📉  No Live Metrics",         "Clinic managers have zero real-time insight into daily doctor workload, patient age demographics, diagnostic patterns, or revenue indicators."),
]

card_w = Inches(3.0)
for i, (title, desc) in enumerate(problems):
    lx = Inches(0.35 + i * 3.24)
    add_card(s, lx, Inches(2.85), card_w, Inches(3.75),
             bg=CARD_BG, border=[BLUE, TEAL, PURPLE, PINK][i])
    add_text(s, title, lx + Inches(0.15), Inches(3.0), card_w - Inches(0.3), Inches(0.45),
             font_size=13, bold=True, color=[BLUE, TEAL, PURPLE, PINK][i])
    add_text(s, desc,  lx + Inches(0.15), Inches(3.52), card_w - Inches(0.3), Inches(2.85),
             font_size=11, color=MUTED)

slide_number(s, 3)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — OBJECTIVES & SCOPE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Objectives & Scope", "Core goals and boundaries of the project", accent=BLUE)

objectives = [
    "Integrated Operations – Bridge gaps between patient registration, doctor scheduling, appointment logging, and invoicing in one synchronized platform.",
    "Automated Invoicing – Business rule engine that auto-triggers a ₹150 billing record whenever a doctor marks a consultation COMPLETED.",
    "Persistent Data Storage – Normalized schema via Spring Data JPA and H2 file-based storage; zero database setup, data preserved across restarts.",
    "Premium SPA Dashboard – Modern UI with real-time KPI metrics, age & revenue charts, and glassmorphic dark/light theming.",
    "Clerical Efficiency – Enable ward managers, billing clerks, and medical officers to quickly search, schedule, update, and audit clinical records.",
]

add_card(s, Inches(0.3), Inches(1.7), Inches(8.0), Inches(5.5), bg=CARD_BG, border=BLUE)
card_title(s, "Project Objectives", Inches(0.5), Inches(1.85), Inches(7.5), accent=BLUE, fs=14)
for i, obj in enumerate(objectives):
    add_rect(s, Inches(0.5), Inches(2.35 + i*0.88), Inches(0.3), Inches(0.3), fill_color=BLUE)
    add_text(s, str(i+1), Inches(0.5), Inches(2.35 + i*0.88), Inches(0.3), Inches(0.3),
             font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s, obj, Inches(0.9), Inches(2.3 + i*0.88), Inches(7.2), Inches(0.8),
             font_size=11, color=LIGHT)

# Scope box (right)
scope = [
    "Localized hospitals and multi-department clinics",
    "Full lifecycle: registration → scheduling → billing",
    "End-to-end: POM config → JPA → REST → SPA UI",
    "Intra-clinic network deployment via H2 AUTO_SERVER",
    "Single-tenant; no multi-hospital partition needed",
]
add_card(s, Inches(8.55), Inches(1.7), Inches(4.45), Inches(3.0), bg=CARD_BG, border=TEAL)
card_title(s, "Project Scope", Inches(8.75), Inches(1.85), Inches(4.0), accent=TEAL, fs=14)
add_multi_para(s, scope, Inches(8.75), Inches(2.35), Inches(4.1), Inches(2.2),
               font_size=11, color=LIGHT, bullet_char="✦ ")

lims = ["Single-tenant only", "H2 → PostgreSQL for multi-hospital", "No payment gateway (mocked)", "Auth is future scope"]
add_card(s, Inches(8.55), Inches(4.85), Inches(4.45), Inches(2.2), bg=CARD_BG, border=PINK)
card_title(s, "Known Limitations", Inches(8.75), Inches(5.0), Inches(4.0), accent=PINK, fs=14)
add_multi_para(s, lims, Inches(8.75), Inches(5.45), Inches(4.1), Inches(1.45),
               font_size=11, color=LIGHT, bullet_char="⚠ ")

slide_number(s, 4)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — THEORETICAL BACKGROUND
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Theoretical Background", "Core frameworks and technologies", accent=PURPLE)

theories = [
    ("☕  Java 21", BLUE,
     ["Modern LTS JDK with enhanced performance", "Records, sealed classes, virtual threads", "Fully supported by Spring Boot 3.x"]),
    ("🍃  Spring Boot", TEAL,
     ["Auto-configuration of Spring beans", "Embedded Tomcat — no external server needed", "Starter dependency management via POM"]),
    ("🗃️  Spring Data JPA", PURPLE,
     ["JpaRepository auto-generates CRUD SQL", "@Entity, @ManyToOne ORM annotations", "Hibernate ORM as reference implementation"]),
    ("💾  H2 Database", PINK,
     ["Lightweight Java-native relational DB", "File mode: jdbc:h2:file:./hospital_db", "AUTO_SERVER for external tool access"]),
    ("🌐  SPA Architecture", BLUE,
     ["Static HTML + JavaScript fetch() calls", "Dynamic JSON rendering without page reload", "REST endpoints serve raw JSON data"]),
    ("📊  Chart.js", TEAL,
     ["Doughnut & bar charts for analytics", "Age distribution, revenue breakdown", "Rendered inside the SPA dashboard"]),
]

cw = Inches(4.15)
ch = Inches(2.4)
for i, (title, clr, pts) in enumerate(theories):
    col = i % 3
    row = i // 3
    lx = Inches(0.3 + col * 4.35)
    ty = Inches(1.7  + row * 2.6)
    add_card(s, lx, ty, cw, ch, bg=CARD_BG, border=clr)
    add_text(s, title, lx+Inches(0.15), ty+Inches(0.08), cw-Inches(0.3), Inches(0.42),
             font_size=13, bold=True, color=clr)
    add_multi_para(s, pts, lx+Inches(0.15), ty+Inches(0.58), cw-Inches(0.3), ch-Inches(0.75),
                   font_size=10.5, color=LIGHT, bullet_char="• ")

slide_number(s, 5)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — TECHNOLOGY STACK
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Technology Stack", "Full end-to-end stack used in this project", accent=TEAL)

stack = [
    ("BACKEND",   BLUE,   ["Java 21", "Spring Boot 3.4.1", "Spring Web (REST)", "Bean Validation", "Hibernate ORM", "Apache Maven"]),
    ("DATABASE",  TEAL,   ["H2 Relational DB", "File-based persistence", "AUTO_SERVER mode", "DDL auto-update", "H2 Web Console", "JDBC SQL"]),
    ("FRONTEND",  PURPLE, ["HTML5", "Vanilla CSS3", "JavaScript (ES6+)", "Fetch API (AJAX)", "Chart.js 4.x", "Glassmorphism UI"]),
    ("DEVOPS",    PINK,   ["Git + GitHub", "Docker + Compose", "Render (Cloud)", "Maven Wrapper", "JAR packaging", "render.yaml"]),
]

for i, (layer, clr, items) in enumerate(stack):
    lx = Inches(0.3 + i * 3.2)
    add_card(s, lx, Inches(1.7), Inches(3.0), Inches(5.5), bg=CARD_BG, border=clr)
    add_rect(s, lx, Inches(1.7), Inches(3.0), Inches(0.5), fill_color=clr)
    add_text(s, layer, lx, Inches(1.7), Inches(3.0), Inches(0.5),
             font_size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        ty = Inches(2.35 + j * 0.77)
        add_rect(s, lx + Inches(0.25), ty + Inches(0.06), Inches(0.25), Inches(0.25), fill_color=clr)
        add_text(s, item, lx + Inches(0.6), ty, Inches(2.25), Inches(0.38),
                 font_size=11.5, color=LIGHT)

slide_number(s, 6)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — SYSTEM ARCHITECTURE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "System Architecture", "Three-layer REST architecture with SPA frontend", accent=BLUE)

layers = [
    ("🌐  SPA Frontend", "HTML5 + Vanilla CSS3 + JavaScript",           BLUE),
    ("REST / JSON API", "AJAX fetch() calls over HTTP",                  MUTED),
    ("⚙️  Spring Boot Controllers", "@RestController — Route & Validate", TEAL),
    ("🔗  Spring Data JPA", "@Repository — ORM / Hibernate",            PURPLE),
    ("💾  H2 File Database", "hospital_db.mv.db — Persistent Storage",  PINK),
]

lx = Inches(1.2)
for i, (label, sub, clr) in enumerate(layers):
    ty = Inches(1.75 + i * 1.08)
    if clr == MUTED:
        add_text(s, "⬇ " + label + " ⬇", Inches(0), ty - Inches(0.1),
                 Inches(13.33), Inches(0.5),
                 font_size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    else:
        add_card(s, lx, ty, Inches(6.2), Inches(0.78), bg=CARD_BG, border=clr)
        add_rect(s, lx, ty, Inches(0.12), Inches(0.78), fill_color=clr)
        add_text(s, label, lx + Inches(0.25), ty + Inches(0.06),
                 Inches(5.7), Inches(0.36), font_size=13, bold=True, color=clr)
        add_text(s, sub,   lx + Inches(0.25), ty + Inches(0.42),
                 Inches(5.7), Inches(0.3),  font_size=11, color=MUTED)

# Right side — explanations
notes = [
    (BLUE,   "Frontend Layer",    "SPA serves from /static. JS fetch() hits /api/* endpoints and renders JSON into the DOM dynamically."),
    (TEAL,   "Controller Layer",  "@RestController classes map HTTP verbs to methods. @Valid ensures Bean Validation before processing."),
    (PURPLE, "Repository Layer",  "JpaRepository auto-generates all CRUD SQL. No boilerplate needed — method naming drives queries."),
    (PINK,   "Database Layer",    "H2 persists data to hospital_db.mv.db in the working directory. Zero installation required."),
]
for i, (clr, title, desc) in enumerate(notes):
    ty = Inches(1.75 + i * 1.4)
    add_card(s, Inches(7.9), ty, Inches(5.1), Inches(1.1), bg=CARD_BG, border=clr)
    add_text(s, title, Inches(8.1), ty + Inches(0.08), Inches(4.8), Inches(0.35),
             font_size=12, bold=True, color=clr)
    add_text(s, desc,  Inches(8.1), ty + Inches(0.42), Inches(4.8), Inches(0.6),
             font_size=10.5, color=MUTED)

slide_number(s, 7)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — PATIENT & DOCTOR MANAGEMENT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Patient & Doctor Management", "Core entity modules — Module 1 & Module 2", accent=TEAL)

# Patient card (left)
add_card(s, Inches(0.3), Inches(1.7), Inches(6.1), Inches(5.5), bg=CARD_BG, border=BLUE)
card_title(s, "👤  Patient Module", Inches(0.5), Inches(1.85), Inches(5.8), accent=BLUE, fs=15)
pat_feats = [
    "Register patients: name, DOB, gender, email, phone, address",
    "Medical history stored as free-text per patient",
    "Real-time search and filter in the SPA table",
    "Full CRUD — create, view, update, delete records",
    "5 sample patients seeded on startup",
    "Input validated: @NotBlank, @Email annotations",
]
add_multi_para(s, pat_feats, Inches(0.5), Inches(2.4), Inches(5.6), Inches(3.4),
               font_size=11.5, color=LIGHT, bullet_char="◆ ")

# Sample patients table
cols_h = ["Name",           "Gender", "Condition"]
rows   = [
    ["James Smith",   "Male",   "Hypertension, Asthma"],
    ["Emily Davis",   "Female", "Penicillin allergy"],
    ["Robert Johnson","Male",   "Type 2 Diabetes"],
    ["Linda Taylor",  "Female", "General checkup"],
    ["Michael Brown", "Male",   "Seasonal allergies"],
]
tx = Inches(0.5)
ty = Inches(5.9)
for ci, ch in enumerate(cols_h):
    add_rect(s, tx + Inches(ci*1.9), ty, Inches(1.85), Inches(0.32), fill_color=BLUE)
    add_text(s, ch, tx + Inches(ci*1.9), ty, Inches(1.85), Inches(0.32),
             font_size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
for ri, row in enumerate(rows):
    for ci, cell in enumerate(row):
        bg = CARD_BG if ri % 2 == 0 else RGBColor(0x1a, 0x25, 0x45)
        add_rect(s, tx + Inches(ci*1.9), ty + Inches(0.32 + ri*0.27),
                 Inches(1.85), Inches(0.27), fill_color=bg)
        add_text(s, cell, tx + Inches(ci*1.9) + Inches(0.05),
                 ty + Inches(0.32 + ri*0.27), Inches(1.75), Inches(0.27),
                 font_size=8.5, color=LIGHT)

# Doctor card (right)
add_card(s, Inches(6.75), Inches(1.7), Inches(6.25), Inches(5.5), bg=CARD_BG, border=TEAL)
card_title(s, "👨‍⚕️  Doctor Module", Inches(6.95), Inches(1.85), Inches(5.9), accent=TEAL, fs=15)
doc_feats = [
    "Manage doctor profiles: name, specialization, room, availability schedule",
    "5 pre-seeded doctors across 5 specializations",
    "Availability displayed as slot text (e.g., Mon–Wed: 9 AM – 4 PM)",
    "Full CRUD — create, view, update, delete doctor records",
]
add_multi_para(s, doc_feats, Inches(6.95), Inches(2.4), Inches(5.9), Inches(1.8),
               font_size=11.5, color=LIGHT, bullet_char="◆ ")

doctors = [
    ("Dr. Sarah Connor",  "Cardiology",        "A-101"),
    ("Dr. John Watson",   "Pediatrics",        "B-205"),
    ("Dr. Gregory House", "Neurology",         "C-310"),
    ("Dr. Meredith Grey", "General Medicine",  "A-105"),
    ("Dr. Stephen Strange","Orthopedics",      "D-404"),
]
cols_dh = ["Doctor",          "Specialization",    "Room"]
dx = Inches(6.95)
dy = Inches(4.35)
for ci, ch in enumerate(cols_dh):
    add_rect(s, dx + Inches(ci*2.0), dy, Inches(1.95), Inches(0.32), fill_color=TEAL)
    add_text(s, ch, dx + Inches(ci*2.0), dy, Inches(1.95), Inches(0.32),
             font_size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
for ri, row in enumerate(doctors):
    for ci, cell in enumerate(row):
        bg = CARD_BG if ri % 2 == 0 else RGBColor(0x1a, 0x25, 0x45)
        add_rect(s, dx + Inches(ci*2.0), dy + Inches(0.32 + ri*0.27),
                 Inches(1.95), Inches(0.27), fill_color=bg)
        add_text(s, cell, dx + Inches(ci*2.0) + Inches(0.05),
                 dy + Inches(0.32 + ri*0.27), Inches(1.85), Inches(0.27),
                 font_size=8.5, color=LIGHT)

slide_number(s, 8)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 — APPOINTMENT & BILLING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Appointment & Billing System", "Smart scheduling with auto-invoice generation", accent=PURPLE)

# Appointment
add_card(s, Inches(0.3), Inches(1.7), Inches(6.1), Inches(5.5), bg=CARD_BG, border=PURPLE)
card_title(s, "📅  Appointment Module", Inches(0.5), Inches(1.85), Inches(5.8), accent=PURPLE, fs=15)
apt = [
    "Book appointments by linking Patient ↔ Doctor with date/time",
    "Status lifecycle: PENDING  →  COMPLETED",
    "Add diagnosis and prescription notes post-consultation",
    "Completing an appointment auto-triggers a $150 billing record",
    "5 sample appointments seeded: 3 completed, 2 pending",
    "Full CRUD: create, list, update, delete appointments",
]
add_multi_para(s, apt, Inches(0.5), Inches(2.4), Inches(5.7), Inches(2.8),
               font_size=11.5, color=LIGHT, bullet_char="◆ ")

# Trigger flow diagram
add_rect(s, Inches(0.5), Inches(5.25), Inches(5.7), Inches(0.04), fill_color=PURPLE)
add_text(s, "⚡  Auto-Billing Trigger Flow", Inches(0.5), Inches(5.35), Inches(5.7), Inches(0.35),
         font_size=12, bold=True, color=PURPLE)
flow_boxes = ["PENDING Appt.", "PUT /complete", "Status → COMPLETED", "Billing auto-created ($150)"]
for i, fb in enumerate(flow_boxes):
    bx = Inches(0.5 + i * 1.42)
    add_rect(s, bx, Inches(5.8), Inches(1.3), Inches(0.5),
             fill_color=[PURPLE, BLUE, TEAL, PINK][i], line_color=None)
    add_text(s, fb, bx, Inches(5.8), Inches(1.3), Inches(0.5),
             font_size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, "▶", Inches(1.8 + i*1.42), Inches(5.87), Inches(0.12), Inches(0.36),
                 font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Billing
add_card(s, Inches(6.75), Inches(1.7), Inches(6.25), Inches(5.5), bg=CARD_BG, border=PINK)
card_title(s, "💳  Billing Module", Inches(6.95), Inches(1.85), Inches(5.9), accent=PINK, fs=15)
bil = [
    "Auto-generated $150 invoice per completed appointment",
    "Status tracking: PAID / UNPAID per record",
    "Insurance provider field stored per invoice",
    "Manual billing also supported for lab work & utilities",
    "Filter bills by patient, status, or date",
]
add_multi_para(s, bil, Inches(6.95), Inches(2.4), Inches(5.9), Inches(2.4),
               font_size=11.5, color=LIGHT, bullet_char="◆ ")

bill_rows = [
    ("James Smith",    "$150.00", "PAID",   "Blue Cross"),
    ("Emily Davis",    "$150.00", "UNPAID", "None"),
    ("Robert Johnson", "$150.00", "PAID",   "Aetna"),
    ("James Smith",    "$75.00",  "UNPAID", "Blue Cross"),
]
bcols = ["Patient", "Amount", "Status", "Insurance"]
bx = Inches(6.95)
by = Inches(5.1)
for ci, bch in enumerate(bcols):
    add_rect(s, bx + Inches(ci*1.5), by, Inches(1.45), Inches(0.3), fill_color=PINK)
    add_text(s, bch, bx + Inches(ci*1.5), by, Inches(1.45), Inches(0.3),
             font_size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
for ri, brow in enumerate(bill_rows):
    for ci, cell in enumerate(brow):
        bg = CARD_BG if ri % 2 == 0 else RGBColor(0x1a, 0x25, 0x45)
        clr = TEAL if cell == "PAID" else (PINK if cell == "UNPAID" else LIGHT)
        add_rect(s, bx+Inches(ci*1.5), by+Inches(0.3+ri*0.27), Inches(1.45), Inches(0.27), fill_color=bg)
        add_text(s, cell, bx+Inches(ci*1.5)+Inches(0.04), by+Inches(0.3+ri*0.27),
                 Inches(1.37), Inches(0.27), font_size=8.5, color=clr)

slide_number(s, 9)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — REST API DESIGN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "REST API Design", "Complete endpoint reference for all modules", accent=TEAL)

api_rows = [
    ("GET",    "/api/patients",                  "Retrieve all patient records",               TEAL),
    ("POST",   "/api/patients",                  "Register a new patient",                     BLUE),
    ("PUT",    "/api/patients/{id}",             "Update patient information",                  PURPLE),
    ("DELETE", "/api/patients/{id}",             "Delete a patient record",                     PINK),
    ("GET",    "/api/doctors",                   "List all doctors and availability",           TEAL),
    ("POST",   "/api/doctors",                   "Create a new doctor profile",                 BLUE),
    ("GET",    "/api/appointments",              "Get all appointment records",                 TEAL),
    ("POST",   "/api/appointments",              "Book a new appointment",                      BLUE),
    ("PUT",    "/api/appointments/{id}/complete","Mark as COMPLETED → auto-triggers billing",  PURPLE),
    ("GET",    "/api/billing",                   "Retrieve all billing records",                TEAL),
    ("PUT",    "/api/billing/{id}/pay",          "Mark an invoice as PAID",                    PURPLE),
    ("GET",    "/api/stats",                     "Live KPI metrics for dashboard widgets",      TEAL),
]

header_cols = ["Method", "Endpoint", "Description"]
col_widths  = [Inches(1.3), Inches(4.8), Inches(6.7)]
tx = Inches(0.3)
ty = Inches(1.7)

for ci, (hcol, cw) in enumerate(zip(header_cols, col_widths)):
    lx = tx + sum(col_widths[:ci])
    add_rect(s, lx, ty, cw - Inches(0.04), Inches(0.38), fill_color=TEAL)
    add_text(s, hcol, lx + Inches(0.08), ty, cw, Inches(0.38),
             font_size=11, bold=True, color=DARK_BG)

for ri, (method, endpoint, desc, clr) in enumerate(api_rows):
    row_ty = ty + Inches(0.38 + ri * 0.42)
    bg = CARD_BG if ri % 2 == 0 else RGBColor(0x14, 0x1f, 0x3a)
    for ci, (val, cw) in enumerate(zip([method, endpoint, desc], col_widths)):
        lx = tx + sum(col_widths[:ci])
        add_rect(s, lx, row_ty, cw - Inches(0.04), Inches(0.40), fill_color=bg)
        txt_clr = clr if ci == 0 else (BLUE if ci == 1 else LIGHT)
        add_text(s, val, lx + Inches(0.08), row_ty + Pt(1),
                 cw - Inches(0.12), Inches(0.36),
                 font_size=10.5 if ci > 0 else 10,
                 bold=(ci == 0), color=txt_clr)

slide_number(s, 10)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 — SPA DASHBOARD & UI
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "SPA Dashboard & UI Design", "Premium single-page interface with real-time analytics", accent=PINK)

ui_feats = [
    ("🌓  Dark / Light Mode",   "CSS custom properties (variables) allow seamless theme switching with no page reload."),
    ("📊  Live KPI Cards",      "Patients, Doctors, Appointments, Revenue counters updated in real-time via /api/stats."),
    ("📈  Chart.js Analytics",  "Bar chart for patient age distribution; doughnut chart for billing status breakdown."),
    ("🔍  Real-time Search",    "Instant filter across all data tables — patient name, doctor specialization, bill status."),
    ("📋  CRUD Forms",          "All entity forms embedded in-page; no page transitions. Smooth SPA experience."),
    ("💎  Glassmorphism",       "Frosted glass effect cards with backdrop-filter, subtle borders and gradient highlights."),
]
for i, (title, desc) in enumerate(ui_feats):
    col = i % 2
    row = i // 2
    lx = Inches(0.3 + col * 6.5)
    ty = Inches(1.7 + row * 1.87)
    add_card(s, lx, ty, Inches(6.1), Inches(1.72), bg=CARD_BG, border=[BLUE,TEAL,PURPLE,PINK,BLUE,TEAL][i])
    add_text(s, title, lx+Inches(0.15), ty+Inches(0.1), Inches(5.8), Inches(0.45),
             font_size=13, bold=True, color=[BLUE,TEAL,PURPLE,PINK,BLUE,TEAL][i])
    add_text(s, desc, lx+Inches(0.15), ty+Inches(0.58), Inches(5.8), Inches(1.0),
             font_size=11, color=MUTED)

# KPI stat blocks
stats = [("Patients","5",BLUE), ("Doctors","5",TEAL), ("Appointments","5",PURPLE), ("Revenue","$525",PINK)]
sw = Inches(3.05)
for i, (label, val, clr) in enumerate(stats):
    lx = Inches(0.3 + i * 3.26)
    add_card(s, lx, Inches(7.0), sw, Inches(0.42), bg=CARD_BG, border=clr)
    add_text(s, f"{val}  {label}", lx, Inches(7.0), sw, Inches(0.42),
             font_size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)

slide_number(s, 11)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 — DATABASE DESIGN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Database Design — JPA & H2", "Entity mapping and relational schema", accent=PURPLE)

entities = [
    ("PATIENT",     BLUE,   ["id  BIGINT (PK)", "name  VARCHAR", "date_of_birth  VARCHAR", "gender  VARCHAR", "email  VARCHAR", "phone  VARCHAR", "address  VARCHAR", "medical_history  TEXT"]),
    ("DOCTOR",      TEAL,   ["id  BIGINT (PK)", "name  VARCHAR", "specialization  VARCHAR", "email  VARCHAR", "phone  VARCHAR", "availability  VARCHAR", "room_number  VARCHAR"]),
    ("APPOINTMENT", PURPLE, ["id  BIGINT (PK)", "patient_id  FK → PATIENT", "doctor_id  FK → DOCTOR", "date_time  VARCHAR", "status  VARCHAR", "symptoms  VARCHAR", "diagnosis  TEXT", "prescription  TEXT"]),
    ("BILLING",     PINK,   ["id  BIGINT (PK)", "patient_id  FK → PATIENT", "appointment_id  FK → APPOINTMENT", "amount  DOUBLE", "status  VARCHAR", "bill_date  VARCHAR", "insurance_provider  VARCHAR"]),
]

ew = Inches(3.0)
for i, (name, clr, fields) in enumerate(entities):
    lx = Inches(0.3 + i * 3.25)
    eh = Inches(0.38 + len(fields) * 0.37)
    add_card(s, lx, Inches(1.7), ew, Inches(5.5), bg=CARD_BG, border=clr)
    add_rect(s, lx, Inches(1.7), ew, Inches(0.45), fill_color=clr)
    add_text(s, name, lx, Inches(1.7), ew, Inches(0.45),
             font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    for fi, field in enumerate(fields):
        ty = Inches(2.22 + fi * 0.37)
        bg = RGBColor(0x1a,0x25,0x45) if fi%2==1 else CARD_BG
        add_rect(s, lx+Inches(0.04), ty, ew-Inches(0.08), Inches(0.34), fill_color=bg)
        fc = clr if "PK" in field or "FK" in field else LIGHT
        add_text(s, field, lx+Inches(0.12), ty+Pt(1), ew-Inches(0.2), Inches(0.32),
                 font_size=9.5, color=fc, bold=("PK" in field or "FK" in field))

# Relationships row
add_rect(s, Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.06), fill_color=PURPLE)
rels = ["Patient  1:N  Appointment", "Doctor  1:N  Appointment", "Appointment  1:1  Billing", "Patient  1:N  Billing"]
for i, rel in enumerate(rels):
    clr = [BLUE, TEAL, PURPLE, PINK][i]
    add_rect(s, Inches(0.3 + i*3.25), Inches(7.25), Inches(3.1), Inches(0.22), fill_color=CARD_BG, line_color=clr)
    add_text(s, rel, Inches(0.3+i*3.25), Inches(7.25), Inches(3.1), Inches(0.22),
             font_size=9, bold=True, color=clr, align=PP_ALIGN.CENTER)

slide_number(s, 12)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13 — DEPLOYMENT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Deployment Strategy", "Local, Docker, and Cloud deployment options", accent=TEAL)

deploy_options = [
    ("🖥️  Local JAR Run", BLUE,
     "java -jar target/management-0.0.1-SNAPSHOT.jar",
     ["Zero setup — run directly on any Java 21 JVM", "H2 file persists data to ./hospital_db", "Access at http://localhost:8080", "Ideal for development and demos"]),
    ("🐳  Docker Compose", TEAL,
     "docker compose up --build -d",
     ["Multi-stage Dockerfile for optimized image size", "Volume mount for H2 persistence across restarts", "Non-root spring user for security", "Port 8080 mapped to host"]),
    ("☁️  Render (Cloud)", PURPLE,
     "render.yaml auto-detected on push",
     ["Connected via GitHub repo webhook", "Build: ./mvnw clean package -DskipTests", "Start: java -Dspring.profiles.active=prod -jar ...", "Free tier — public HTTPS URL assigned"]),
]

for i, (title, clr, cmd, pts) in enumerate(deploy_options):
    lx = Inches(0.3 + i * 4.35)
    add_card(s, lx, Inches(1.7), Inches(4.15), Inches(5.5), bg=CARD_BG, border=clr)
    add_rect(s, lx, Inches(1.7), Inches(4.15), Inches(0.5), fill_color=clr)
    add_text(s, title, lx+Inches(0.1), Inches(1.7), Inches(4.0), Inches(0.5),
             font_size=13, bold=True, color=DARK_BG)
    add_card(s, lx+Inches(0.1), Inches(2.32), Inches(3.9), Inches(0.48),
             bg=RGBColor(0x08,0x0d,0x18), border=clr)
    add_text(s, cmd, lx+Inches(0.18), Inches(2.32), Inches(3.78), Inches(0.48),
             font_size=9, color=TEAL, font_name="Courier New")
    add_multi_para(s, pts, lx+Inches(0.15), Inches(2.9), Inches(3.85), Inches(2.5),
                   font_size=11, color=LIGHT, bullet_char="✦ ")

# Build success banner
add_rect(s, Inches(0.3), Inches(7.08), Inches(12.7), Inches(0.36), fill_color=RGBColor(0x00,0x2d,0x22))
add_rect(s, Inches(0.3), Inches(7.08), Inches(0.1), Inches(0.36), fill_color=TEAL)
add_text(s, "✅  BUILD SUCCESS — JAR verified: ./mvnw clean package -DskipTests  → target/management-0.0.1-SNAPSHOT.jar",
         Inches(0.5), Inches(7.1), Inches(12.5), Inches(0.32), font_size=11, bold=True, color=TEAL)

slide_number(s, 13)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14 — REAL-WORLD EVALUATION & LIMITATIONS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)
add_header_bar(s, "Real-World Evaluation & Future Scope", "Clinical validation and planned enhancements", accent=PINK)

# Evaluation (left)
add_card(s, Inches(0.3), Inches(1.7), Inches(6.1), Inches(5.5), bg=CARD_BG, border=BLUE)
card_title(s, "🏥  Clinical Evaluation", Inches(0.5), Inches(1.85), Inches(5.8), accent=BLUE, fs=15)

add_text(s, "Kalyan Hospital, Raipur", Inches(0.5), Inches(2.4), Inches(5.6), Inches(0.38),
         font_size=13, bold=True, color=TEAL)
add_text(s, "Dr. Sameer Sen — Senior Medical Officer",
         Inches(0.5), Inches(2.78), Inches(5.6), Inches(0.3), font_size=11, color=MUTED)
kh = ["Operational appointment monitoring", "Doctor availability tracking", "Real-time patient record retrieval"]
add_multi_para(s, kh, Inches(0.5), Inches(3.08), Inches(5.6), Inches(0.9),
               font_size=10.5, color=LIGHT, bullet_char="◆ ")

add_rect(s, Inches(0.5), Inches(4.1), Inches(5.5), Inches(0.02), fill_color=PURPLE)

add_text(s, "BSP Sector 7 General Hospital, Raipur", Inches(0.5), Inches(4.2),
         Inches(5.6), Inches(0.38), font_size=13, bold=True, color=PURPLE)
add_text(s, "Sister Mini Kurian — Chief Nursing Superintendent",
         Inches(0.5), Inches(4.58), Inches(5.6), Inches(0.3), font_size=11, color=MUTED)
bh = ["Patient registration and ward allocations", "Billing and invoice audits", "Ward nurse shift logging"]
add_multi_para(s, bh, Inches(0.5), Inches(4.88), Inches(5.6), Inches(0.9),
               font_size=10.5, color=LIGHT, bullet_char="◆ ")

# Key results
results = [("~50ms", "Avg API Response"), ("0", "DB Setup Steps"), ("2", "Hospitals Evaluated"), ("100%", "Build Pass Rate")]
for i, (val, lbl) in enumerate(results):
    rx = Inches(0.5 + i * 1.38)
    add_rect(s, rx, Inches(5.95), Inches(1.28), Inches(0.8), fill_color=RGBColor(0x0d,0x18,0x33), line_color=[BLUE,TEAL,PURPLE,PINK][i])
    add_text(s, val, rx, Inches(5.97), Inches(1.28), Inches(0.45), font_size=16, bold=True, color=[BLUE,TEAL,PURPLE,PINK][i], align=PP_ALIGN.CENTER)
    add_text(s, lbl, rx, Inches(6.42), Inches(1.28), Inches(0.3), font_size=8, color=MUTED, align=PP_ALIGN.CENTER)

# Future scope (right)
add_card(s, Inches(6.75), Inches(1.7), Inches(6.25), Inches(5.5), bg=CARD_BG, border=TEAL)
card_title(s, "🔮  Future Scope", Inches(6.95), Inches(1.85), Inches(5.9), accent=TEAL, fs=15)
future = [
    "Role-based authentication — Admin, Doctor, Nurse portals",
    "PostgreSQL / MySQL migration for multi-hospital support",
    "Real payment gateway integration (Razorpay / Stripe)",
    "Mobile-responsive Progressive Web App (PWA) with offline mode",
    "WhatsApp / SMS appointment reminders via Twilio API",
    "Electronic Health Records (EHR) module with PDF export",
    "AI-based diagnosis suggestion engine",
    "Microservices split for independent scaling of modules",
]
add_multi_para(s, future, Inches(6.95), Inches(2.4), Inches(5.9), Inches(4.2),
               font_size=11.5, color=LIGHT, bullet_char="◆ ")

slide_number(s, 14)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15 — CONCLUSION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank_layout)
add_bg(s)

# Gradient top stripe
add_rect(s, Inches(0), Inches(0), Inches(13.33), Inches(0.12), fill_color=TEAL)
add_rect(s, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3),  fill_color=CARD_BG)

add_text(s, "Conclusion", Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.75),
         font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

summary = [
    "Full-stack Spring Boot Hospital Management Platform — 0 to Production",
    "Automated billing engine eliminates manual invoice creation on every consultation",
    "Premium SPA dashboard with real-time analytics, charts, and dark/light theming",
    "Deployed and validated at 2 real hospitals — Kalyan Hospital & BSP Sector 7",
    "Containerized with Docker & live on Render cloud — shareable public URL",
    "Zero database installation — H2 file-based engine with AUTO_SERVER mode",
    "All REST APIs follow standard HTTP semantics (GET/POST/PUT/DELETE)",
    "Built with Java 21, Spring Boot 3.4, H2, Maven, Chart.js, Vanilla JS SPA",
]
add_card(s, Inches(0.3), Inches(1.08), Inches(8.0), Inches(6.1), bg=CARD_BG, border=TEAL)
card_title(s, "✅  What Was Achieved", Inches(0.5), Inches(1.22), Inches(7.7), accent=TEAL, fs=14)
add_multi_para(s, summary, Inches(0.5), Inches(1.7), Inches(7.6), Inches(5.0),
               font_size=11.5, color=LIGHT, bullet_char="◆ ")

# Thank you box
add_card(s, Inches(8.6), Inches(1.08), Inches(4.4), Inches(6.1), bg=CARD_BG, border=BLUE)
add_text(s, "🙏", Inches(8.6), Inches(1.4), Inches(4.4), Inches(0.7), font_size=36, align=PP_ALIGN.CENTER)
add_text(s, "Thank You!", Inches(8.6), Inches(2.15), Inches(4.4), Inches(0.7),
         font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(s, "Questions & Discussion", Inches(8.6), Inches(2.9), Inches(4.4), Inches(0.38),
         font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s, Inches(9.2), Inches(3.42), Inches(3.2), Inches(0.04), fill_color=BLUE)

add_text(s, "S Aditya", Inches(8.6), Inches(3.55), Inches(4.4), Inches(0.5),
         font_size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, "M.Sc. Computer Science", Inches(8.6), Inches(4.05), Inches(4.4), Inches(0.35),
         font_size=12, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, "ICFAI University, Raipur", Inches(8.6), Inches(4.4), Inches(4.4), Inches(0.35),
         font_size=12, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, "Enrollment: 22STUCRPN01019", Inches(8.6), Inches(4.75), Inches(4.4), Inches(0.35),
         font_size=11, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, "Session: 2024 – 26", Inches(8.6), Inches(5.1), Inches(4.4), Inches(0.35),
         font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Tech badges
badges = [("Java 21", BLUE), ("Spring Boot", TEAL), ("H2 + JPA", PURPLE), ("SPA + Chart.js", PINK)]
for i, (name, clr) in enumerate(badges):
    bx = Inches(8.7 + (i%2)*2.05)
    by = Inches(5.7 + (i//2)*0.5)
    add_rect(s, bx, by, Inches(1.95), Inches(0.38), fill_color=RGBColor(0x12,0x1a,0x35), line_color=clr)
    add_text(s, name, bx, by, Inches(1.95), Inches(0.38),
             font_size=10, bold=True, color=clr, align=PP_ALIGN.CENTER)

slide_number(s, 15)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "Hospital_Management_System_Presentation.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
